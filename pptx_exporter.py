import io
import logging
import re
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml
from svgelements import (
    SVG,
    Path,
    Shape,
    Color,
    Move,
    Line,
    CubicBezier,
    QuadraticBezier,
    Arc,
    Close,
    Text,
    Image,
    Group,
    Use,
)

try:
    from resvg_py import resvg_py
except ImportError:
    resvg_py = None

    import cairosvg
except ImportError:
    cairosvg = None

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class PPTXNativeExporter:
    @staticmethod

    def _color_to_rgb(color_obj, svg_doc=None, element_values=None, svg_code=None):
        """Convert svgelements Color to python-pptx RGBColor.
        
        Args:
            color_obj: The color object (usually a Color instance)
            svg_doc: The SVG document (optional, for resolving gradients)
            element_values: The raw element values dict (optional, for detecting gradient URLs)
            svg_code: The raw SVG code string (optional, for parsing gradient stops via regex)
        """
        import re
        
        if not color_obj:
            return None
        
        # Check if this is a gradient URL in the raw values
        # This happens when svgelements couldn't resolve the gradient and falls back to black
        if element_values and isinstance(element_values.get('fill'), str):
            fill_str = element_values.get('fill')
            if fill_str.startswith('url(#') and fill_str.endswith(')'):
                # Extract gradient ID: url(#gradientId) -> gradientId
                gradient_id = fill_str[5:-1]
                logger.debug(f"Detected gradient URL: {fill_str} -> ID: {gradient_id}")
                
                # Try to extract first stop color from gradient using regex on SVG code
                if svg_code:
                    try:
                        # Build regex to find the gradient definition and its stops
                        # Matches both linearGradient and radialGradient
                        gradient_pattern = rf'<(linear|radial)Gradient[^>]*id=["\']?{re.escape(gradient_id)}["\']?[^>]*>.*?</(linear|radial)Gradient>'
                        gradient_match = re.search(gradient_pattern, svg_code, re.DOTALL)
                        
                        if gradient_match:
                            gradient_def = gradient_match.group(0)
                            logger.debug(f"Found gradient definition")
                            
                            # Extract stop-color from the first stop
                            # Handle both style="stop-color:#color" and fill="#color" attributes
                            stop_pattern = r'<stop[^>]*style=["\']([^"\']*)'
                            stop_match = re.search(stop_pattern, gradient_def)
                            
                            if stop_match:
                                stop_style = stop_match.group(1)
                                # Extract stop-color from style attribute
                                color_match = re.search(r'stop-color:\s*([^;\s]+)', stop_style)
                                if color_match:
                                    stop_color_str = color_match.group(1).strip()
                                    logger.debug(f"Extracted stop color from gradient: {stop_color_str}")
                                    # Parse the color string (e.g., 'rgb(135, 206, 235)' or '#87CEEB')
                                    try:
                                        stop_color_obj = Color(stop_color_str)
                                        return RGBColor(stop_color_obj.red, stop_color_obj.green, stop_color_obj.blue)
                                    except Exception as parse_err:
                                        logger.debug(f"Could not parse stop color {stop_color_str}: {parse_err}")
                    except Exception as gradient_err:
                        logger.warning(f"Failed to extract gradient color: {gradient_err}")
                
                # Fallback for gradients: use a light gray instead of black
                logger.warning(f"Could not resolve gradient {gradient_id}, using light gray fallback")
                return RGBColor(200, 200, 200)
        
        # Standard color object handling
        if hasattr(color_obj, 'value') and color_obj.value is None:
            return None
        
        try:
            return RGBColor(color_obj.red, color_obj.green, color_obj.blue)
        except Exception as e:
            logger.warning(f"Failed to convert color {color_obj}: {e}")
            return RGBColor(128, 128, 128)

    @staticmethod
    def _apply_transparency(color_format, alpha_int):
        """
        Applies transparency to a ColorFormat object by manipulating XML.
        alpha_int: 0 (transparent) to 255 (opaque).
        """
        if alpha_int is None or alpha_int >= 255:
            return

        try:
            if hasattr(color_format._color, "_xClr"):
                clr_element = color_format._color._xClr
            else:
                return

            # XML alpha is 0-100000 (100% opacity = 100000)
            xml_alpha_val = int((alpha_int / 255.0) * 100000)

            # Namespace map for 'a'
            ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}

            # Remove existing alpha if any
            existing_alpha = clr_element.find("a:alpha", ns)
            if existing_alpha is not None:
                clr_element.remove(existing_alpha)

            # Create new alpha element
            alpha_elem = parse_xml(
                f'<a:alpha xmlns:a="{ns["a"]}" val="{xml_alpha_val}"/>'
            )
            clr_element.append(alpha_elem)

        except Exception as e:
            logger.warning(f"Failed to apply transparency: {e}")

    @staticmethod

    def _is_gradient_fill(element):
        """
        Check if an element's fill is a gradient url() reference.
        
        Returns:
            bool: True if fill is a gradient (url(#...)), False otherwise
        """
        if not hasattr(element, 'fill'):
            return False
        
        fill = element.fill
        
        # Check if fill is a url() reference (e.g., 'url(#gradient-id)')
        # This is the primary way svgelements represents gradient fills
        if isinstance(fill, str) and fill.startswith('url(#') and fill.endswith(')'):
            logger.debug(f"Detected gradient URL: {fill}")
            return True
        
        # Also check raw values dict for url() references
        if hasattr(element, 'values') and isinstance(element.values.get('fill'), str):
            fill_str = element.values.get('fill')
            if fill_str.startswith('url(#') and fill_str.endswith(')'):
                logger.debug(f"Detected gradient URL in raw values: {fill_str}")
                return True
        
        return False


    @staticmethod

    def _build_gradient_svg(element, svg_doc=None, svg_code=None):
        """
        Build a minimal SVG string containing only the element and its gradient definition.
        Extracts the gradient from the original SVG code or tries to infer it.
        
        Args:
            element: The SVG element with gradient fill
            svg_doc: The SVG document (optional, for getting bounding box)
            svg_code: The original SVG code (required, to extract gradient definitions)
        
        Returns:
            str: SVG XML string, or None if reconstruction fails
        """
        try:
            # Get gradient ID from fill attribute
            gradient_id = None
            if hasattr(element, 'values') and isinstance(element.values.get('fill'), str):
                fill_str = element.values.get('fill')
                if fill_str.startswith('url(#') and fill_str.endswith(')'):
                    gradient_id = fill_str[5:-1]  # Extract ID from url(#id)
            
            if not gradient_id:
                logger.warning("Could not extract gradient ID")
                return None
            
            # Extract gradient definition from SVG code
            if svg_code:
                gradient_def_xml = PPTXNativeExporter._extract_gradient_from_svg_code(gradient_id, svg_code)
                if not gradient_def_xml:
                    logger.warning(f"Could not find gradient definition for {gradient_id}")
                    return None
            else:
                logger.warning("SVG code not provided; cannot extract gradient definition")
                return None
            
            # Get element bounding box for viewBox
            try:
                bbox = element.bbox()
                if bbox and len(bbox) >= 4:
                    x_min, y_min, x_max, y_max = bbox[0], bbox[1], bbox[2], bbox[3]
                else:
                    x_min, y_min, x_max, y_max = 0, 0, 100, 100
            except:
                x_min, y_min, x_max, y_max = 0, 0, 100, 100
            
            width = max(1, x_max - x_min)
            height = max(1, y_max - y_min)
            
            # Build minimal SVG: svg > defs > gradient > element
            svg_lines = [
                f'<svg xmlns="http://www.w3.org/2000/svg" ',
                f'xmlns:xlink="http://www.w3.org/1999/xlink" ',
                f'width="{width}" height="{height}" ',
                f'viewBox="{x_min} {y_min} {width} {height}">',
                f'<defs>{gradient_def_xml}</defs>',
            ]
            
            # Reconstruct element as SVG  
            element_xml = PPTXNativeExporter._element_to_svg_xml(element, gradient_id)
            if not element_xml:
                logger.warning("Failed to convert element to SVG XML")
                return None
            
            svg_lines.append(element_xml)
            svg_lines.append('</svg>')
            
            result = "".join(svg_lines)
            logger.debug(f"Generated gradient SVG ({len(result)} bytes)")
            return result
        
        except Exception as e:
            logger.error(f"Failed to build gradient SVG: {e}")
            return None

    @staticmethod
    def _extract_gradient_from_svg_code(gradient_id, svg_code):
        """
        Extract a gradient definition from SVG code using regex.
        
        Args:
            gradient_id: The gradient ID to find (e.g., 'myGrad')
            svg_code: The SVG XML as a string
        
        Returns:
            str: The gradient XML definition, or None if not found
        """
        try:
            # Pattern to match linearGradient or radialGradient with the given ID
            pattern = rf'<(linear|radial)Gradient[^>]*id=["\']?{re.escape(gradient_id)}["\']?[^>]*>.*?</(linear|radial)Gradient>'
            match = re.search(pattern, svg_code, re.DOTALL | re.IGNORECASE)
            
            if match:
                gradient_xml = match.group(0)
                logger.debug(f"Extracted gradient definition: {len(gradient_xml)} bytes")
                return gradient_xml
            else:
                logger.warning(f"Gradient definition not found for ID: {gradient_id}")
                return None
        
        except Exception as e:
            logger.error(f"Failed to extract gradient from SVG code: {e}")
            return None


        try:
            fill = element.fill
            
            # Get element bounding box for viewBox
            try:
                bbox = element.bbox()
                if bbox and len(bbox) >= 4:
                    x_min, y_min, x_max, y_max = bbox[0], bbox[1], bbox[2], bbox[3]
                else:
                    x_min, y_min, x_max, y_max = 0, 0, 100, 100
            except:
                x_min, y_min, x_max, y_max = 0, 0, 100, 100
            
            width = max(1, x_max - x_min)
            height = max(1, y_max - y_min)
            
            # Start SVG container
            svg_lines = [
                f'<svg xmlns="http://www.w3.org/2000/svg" '
                f'width="{width}" height="{height}" '
                f'viewBox="{x_min} {y_min} {width} {height}">'
            ]
            
            # Add gradient definition
            if isinstance(fill, LinearGradient):
                grad_id = "grad_linear"
                gradient_xml = PPTXNativeExporter._build_linear_gradient_xml(fill, grad_id)
                svg_lines.append(f"<defs>{gradient_xml}</defs>")
            elif isinstance(fill, RadialGradient):
                grad_id = "grad_radial"
                gradient_xml = PPTXNativeExporter._build_radial_gradient_xml(fill, grad_id)
                svg_lines.append(f"<defs>{gradient_xml}</defs>")
            else:
                logger.warning("Cannot reconstruct gradient: unknown type")
                return None
            
            # Reconstruct element as SVG path/shape
            # Extract path data or shape attributes
            element_xml = PPTXNativeExporter._element_to_svg_xml(element, grad_id)
            if not element_xml:
                logger.warning("Failed to convert element to SVG XML")
                return None
            
            svg_lines.append(element_xml)
            svg_lines.append("</svg>")
            
            result = "".join(svg_lines)
            logger.debug(f"Generated gradient SVG ({len(result)} bytes)")
            return result
        
        except Exception as e:
            logger.error(f"Failed to build gradient SVG: {e}")
            return None


    @staticmethod
    def _element_to_svg_xml(element, grad_id):
        """
        Convert an SVG element to its XML representation with gradient reference.
        
        Args:
            element: The SVG element
            grad_id: Gradient ID to reference in fill
        
        Returns:
            str: XML string representing the element
        """
        try:
            from svgelements import Rect, Circle, Ellipse, Line as SvgLine, Polygon, Polyline
            
            # Handle different element types
            if isinstance(element, Path):
                # Convert path data
                try:
                    d = element.d()
                except:
                    d = ""
                if not d:
                    logger.warning("Could not extract path data")
                    return None
                return f'<path d="{d}" fill="url(#{grad_id})"/>'
            
            elif isinstance(element, Rect):
                x = element.x if hasattr(element, 'x') else 0
                y = element.y if hasattr(element, 'y') else 0
                width = element.width if hasattr(element, 'width') else 100
                height = element.height if hasattr(element, 'height') else 100
                rx = element.rx if hasattr(element, 'rx') and element.rx else 0
                ry = element.ry if hasattr(element, 'ry') and element.ry else 0
                
                rx_attr = f' rx="{rx}"' if rx else ''
                ry_attr = f' ry="{ry}"' if ry else ''
                
                return f'<rect x="{x}" y="{y}" width="{width}" height="{height}"{rx_attr}{ry_attr} fill="url(#{grad_id})"/>'
            
            elif isinstance(element, Circle):
                cx = element.cx if hasattr(element, 'cx') else 50
                cy = element.cy if hasattr(element, 'cy') else 50
                r = element.r if hasattr(element, 'r') else 50
                return f'<circle cx="{cx}" cy="{cy}" r="{r}" fill="url(#{grad_id})"/>'
            
            elif isinstance(element, Ellipse):
                cx = element.cx if hasattr(element, 'cx') else 50
                cy = element.cy if hasattr(element, 'cy') else 50
                rx = element.rx if hasattr(element, 'rx') else 50
                ry = element.ry if hasattr(element, 'ry') else 30
                return f'<ellipse cx="{cx}" cy="{cy}" rx="{rx}" ry="{ry}" fill="url(#{grad_id})"/>'
            
            else:
                # Fallback: try to render as path
                try:
                    path_obj = Path(element)
                    d = path_obj.d()
                    if d:
                        return f'<path d="{d}" fill="url(#{grad_id})"/>'
                except:
                    pass
                
                logger.warning(f"Could not convert element type {type(element).__name__} to SVG")
                return None
        
        except Exception as e:
            logger.error(f"Failed to convert element to SVG XML: {e}")
            return None

    @staticmethod
    def _rasterize_element_to_png(svg_string, element):
        """
        Convert an SVG string (containing a gradient element) to PNG bytes using resvg-py.
        
        Args:
            svg_string: SVG XML as string
            element: Original SVG element (for reference, to extract dimensions)
        
        Returns:
            bytes: PNG image data, or None if conversion fails
        """
        if not resvg_py:
            logger.error("resvg_py not available; gradient rasterization disabled")
            return None
        
        try:
            # Convert SVG to PNG using resvg-py (Rust-based, no external DLLs needed)
            # svg_to_bytes returns bytes directly
            png_data = resvg_py.svg_to_bytes(svg_string)
            
            if png_data:
                logger.debug(f"Rasterized gradient element to PNG ({len(png_data)} bytes)")
                return png_data
            else:
                logger.warning("resvg-py produced empty PNG")
                return None
        
        except Exception as e:
            logger.error(f"Failed to rasterize gradient element: {e}")
            return None


    @staticmethod
    def _approximate_cubic_bezier(p0, p1, p2, p3, steps=10):
        """
        Approximate a cubic Bezier curve with line segments.
        Uses the cubic Bezier formula: B(t) = (1-t)^3*P0 + 3(1-t)^2*t*P1 + 3(1-t)*t^2*P2 + t^3*P3
        Returns list of (x, y) tuples representing points on the curve.
        
        Args:
            p0, p1, p2, p3: Control points (with .x, .y attributes)
            steps: Number of line segments to approximate with
        """
        points = []
        for i in range(1, steps + 1):
            t = i / steps
            # Calculate (1-t) and t squared/cubed for efficiency
            one_minus_t = 1 - t
            one_minus_t_sq = one_minus_t * one_minus_t
            one_minus_t_cu = one_minus_t_sq * one_minus_t
            t_sq = t * t
            t_cu = t_sq * t
            
            # Apply cubic Bezier formula
            x = (one_minus_t_cu * p0.x + 
                 3 * one_minus_t_sq * t * p1.x + 
                 3 * one_minus_t * t_sq * p2.x + 
                 t_cu * p3.x)
            y = (one_minus_t_cu * p0.y + 
                 3 * one_minus_t_sq * t * p1.y + 
                 3 * one_minus_t * t_sq * p2.y + 
                 t_cu * p3.y)
            
            points.append((x, y))
        return points


    @staticmethod
    def _bezier_to_line_segments(bezier_curve, tx, ty, num_points=10):
        """
        Approximate a Bezier curve with line segments.
        Returns list of (x, y) tuples in PPTX coordinate units.
        """
        segments = []
        for i in range(1, num_points + 1):
            t = i / num_points
            # Use the point() method to evaluate curve at parameter t
            try:
                point = bezier_curve.point(t)
                segments.append((tx(point.x), ty(point.y)))
            except Exception as e:
                logger.debug(f"Error evaluating bezier curve at t={t}: {e}")
        return segments

    @staticmethod
    def _arc_to_line_segments(arc, tx, ty, num_points=15):
        """
        Approximate an Arc with line segments.
        Returns list of (x, y) tuples in PPTX coordinate units.
        """
        segments = []
        try:
            # Try to convert arc to cubic bezier curves
            cubic_curves = arc.as_cubic_curves()
            for cubic in cubic_curves:
                # Approximate each cubic curve with line segments using the direct formula
                approx_points = PPTXNativeExporter._approximate_cubic_bezier(
                    cubic.start, cubic.control1, cubic.control2, cubic.end, steps=num_points
                )
                # Convert to PPTX coordinates and add to segments
                for p in approx_points:
                    segments.append((tx(p[0]), ty(p[1])))
        except Exception as e:
            logger.debug(f"Failed to convert arc to cubic curves: {e}")
            # Fallback: just the endpoint
            segments.append((tx(arc.end.x), ty(arc.end.y)))
        return segments

    @staticmethod

    def _handle_path(element, slide, tx, ty, svg_doc=None, svg_code=None):
        """
        Handle SVG Path/Shape elements and convert to PPTX shapes.
        
        CRITICAL FIX: Batch all vertices into a single list and call add_line_segments ONCE
        with close=False to avoid creating multiple disconnected closed segments.
        
        Root cause: add_line_segments has close=True by default, so calling it multiple times
        (once per segment) resulted in XML like:
          <a:moveTo>...</a:moveTo>
          <a:lnTo>...</a:lnTo>
          <a:close/>
          <a:lnTo>...
        This creates disjointed segments instead of one continuous path.
        """
        # ===== GRADIENT FALLBACK: Attempt rasterization =====
        if PPTXNativeExporter._is_gradient_fill(element):
            logger.info("Element has gradient fill; attempting rasterization to PNG")
            
            try:
                # Build a minimal SVG containing just this element and its gradient
                svg_string = PPTXNativeExporter._build_gradient_svg(element, svg_doc=svg_doc, svg_code=svg_code)
                
                if svg_string:
                    # Rasterize to PNG
                    png_bytes = PPTXNativeExporter._rasterize_element_to_png(svg_string, element)
                    
                    if png_bytes:
                        # Get element bounding box for positioning
                        try:
                            bbox = element.bbox()
                            if bbox and len(bbox) >= 4:
                                left = tx(bbox[0])
                                top = ty(bbox[1])
                                width = Inches((bbox[2] - bbox[0]) * 0.01)  # Convert px to inches
                                height = Inches((bbox[3] - bbox[1]) * 0.01)
                            else:
                                left, top = Inches(0), Inches(0)
                                width, height = Inches(1), Inches(1)
                        except:
                            left, top = Inches(0), Inches(0)
                            width, height = Inches(1), Inches(1)
                        
                        # Insert PNG as image
                        try:
                            picture = slide.shapes.add_picture(
                                BytesIO(png_bytes),
                                left, top,
                                width=width,
                                height=height
                            )
                            logger.info(f"Gradient rasterized and inserted as image at ({left}, {top})")
                            return picture  # Success!
                        except Exception as img_err:
                            logger.error(f"Failed to insert rasterized gradient image: {img_err}")
                else:
                    logger.warning("Could not build gradient SVG string")
            
            except Exception as grad_err:
                logger.warning(f"Gradient rasterization failed, falling back to vector: {grad_err}")
        
        # ===== Standard Vector Rendering (or fallback if rasterization failed) =====


        try:
            path = Path(element)
            path.reify()  # Apply transformation matrix

            if len(path) == 0:
                logger.debug("Path has no segments")
                return

            start_seg = path[0]
            if not isinstance(start_seg, Move):
                start_x, start_y = 0, 0
            else:
                start_x, start_y = start_seg.end.x, start_seg.end.y

            builder = slide.shapes.build_freeform(tx(start_x), ty(start_y))

            iterator = iter(path)
            next(iterator, None)  # Skip first Move

            # CRITICAL: Collect ALL vertices into a single list
            all_vertices = []
            is_path_closed = False
            current_x, current_y = start_x, start_y  # Track current pen position
            sub_path_start_x, sub_path_start_y = start_x, start_y  # Track sub-path start

            for seg in iterator:
                if isinstance(seg, (Move, Line)):
                    all_vertices.append((tx(seg.end.x), ty(seg.end.y)))
                    current_x, current_y = seg.end.x, seg.end.y
                    # New Move segment starts a new sub-path
                    if isinstance(seg, Move):
                        sub_path_start_x, sub_path_start_y = seg.end.x, seg.end.y

                elif isinstance(seg, CubicBezier):
                    # Approximate cubic Bezier with line segments using direct formula
                    approx_points = PPTXNativeExporter._approximate_cubic_bezier(
                        seg.start, seg.control1, seg.control2, seg.end, steps=10
                    )
                    # Convert to PPTX coordinates and add to vertices
                    for p in approx_points:
                        all_vertices.append((tx(p[0]), ty(p[1])))
                    current_x, current_y = seg.end.x, seg.end.y

                elif isinstance(seg, QuadraticBezier):
                    # Convert QuadraticBezier to CubicBezier, then approximate
                    # Quadratic to Cubic conversion:
                    # For quad bezier Q(t) with control points Q0, Q1, Q2,
                    # the equivalent cubic has control points: C0=Q0, C1=Q0+2/3*(Q1-Q0), C2=Q2+2/3*(Q1-Q2), C3=Q2
                    p0 = seg.start
                    p1_quad = seg.control
                    p2 = seg.end

                    # Convert to cubic control points
                    p1_cubic = type('Point', (), {'x': p0.x + (2/3) * (p1_quad.x - p0.x), 'y': p0.y + (2/3) * (p1_quad.y - p0.y)})()
                    p2_cubic = type('Point', (), {'x': p2.x + (2/3) * (p1_quad.x - p2.x), 'y': p2.y + (2/3) * (p1_quad.y - p2.y)})()

                    approx_points = PPTXNativeExporter._approximate_cubic_bezier(
                        p0, p1_cubic, p2_cubic, p2, steps=10
                    )
                    # Convert to PPTX coordinates and add to vertices
                    for p in approx_points:
                        all_vertices.append((tx(p[0]), ty(p[1])))
                    current_x, current_y = seg.end.x, seg.end.y

                elif isinstance(seg, Arc):
                    # Convert Arc to line segments
                    segments = PPTXNativeExporter._arc_to_line_segments(
                        seg, tx, ty, num_points=15
                    )
                    if segments:
                        all_vertices.extend(segments)
                    current_x, current_y = seg.end.x, seg.end.y

                elif isinstance(seg, Close):
                    # ROBUST CLOSURE HANDLING
                    # Check if current pen position differs from sub-path start
                    # Use small epsilon for floating point comparison
                    epsilon = 1e-6
                    if abs(current_x - sub_path_start_x) > epsilon or abs(current_y - sub_path_start_y) > epsilon:
                        # Explicitly add line segment back to start point
                        logger.debug(
                            f"Closing path: current=({current_x}, {current_y}), "
                            f"start=({sub_path_start_x}, {sub_path_start_y})"
                        )
                        all_vertices.append((tx(sub_path_start_x), ty(sub_path_start_y)))
                        current_x, current_y = sub_path_start_x, sub_path_start_y
                    # Mark path as closed
                    is_path_closed = True
                    logger.debug("Path marked as closed")

            if not all_vertices:
                logger.debug("Path has no drawable segments")
                return None

            # CRITICAL FIX: Call add_line_segments ONCE with ALL vertices and close=False
            # This prevents the default close=True from creating multiple disconnected closed segments.
            # The close=True default was causing XML like:
            #   <a:moveTo>...</a:moveTo>
            #   <a:lnTo>...</a:lnTo>
            #   <a:close/>
            # to be emitted after EACH call, resulting in invisible/disjointed paths.
            logger.debug(f"Adding {len(all_vertices)} vertices to path (all at once)")
            builder.add_line_segments(all_vertices, close=False)

            # Now set builder.closed if path should be closed
            if is_path_closed:
                builder.closed = True
                logger.debug("Path marked as closed for rendering")

            shape = builder.convert_to_shape()
            logger.debug("Path shape created successfully")

            # --- Fill Logic with Opacity Support ---
            if element.fill.value is not None:
                shape.fill.solid()
                rgb = PPTXNativeExporter._color_to_rgb(element.fill, svg_doc, element.values, svg_code)
                if rgb:
                    shape.fill.fore_color.rgb = rgb
                    logger.debug(f"Fill applied: {rgb}")
                    # Transparency handling: multiply fill alpha with element opacity
                    fill_alpha = 255  # Default to fully opaque
                    if hasattr(element.fill, "alpha") and element.fill.alpha is not None:
                        fill_alpha = element.fill.alpha
                    # Apply element-level opacity if present
                    opacity_val = element.values.get('opacity')  # Read opacity attribute
                    if opacity_val is not None:
                        try:
                            opacity_float = float(opacity_val)
                            # opacity_float is 0.0-1.0, convert to 0-255 scale
                            opacity_alpha = int(opacity_float * 255)
                            # Multiply with existing fill alpha
                            fill_alpha = int((fill_alpha / 255.0) * opacity_alpha)
                            logger.debug(f"Opacity applied: {opacity_val} -> alpha={fill_alpha}")
                        except (ValueError, TypeError):
                            logger.debug(f"Could not parse opacity: {opacity_val}")
                    if fill_alpha < 255:
                        PPTXNativeExporter._apply_transparency(
                            shape.fill.fore_color, fill_alpha
                        )
            else:
                shape.fill.background()


            # --- Stroke Logic with Opacity Support ---
            if element.stroke.value is not None and element.stroke_width > 0:
                rgb = PPTXNativeExporter._color_to_rgb(element.stroke, svg_doc, element.values, svg_code)
                if rgb:
                    shape.line.color.rgb = rgb
                    logger.debug(f"Stroke applied: {rgb}, Width: {shape.line.width}")
                    # Transparency handling for stroke
                    stroke_alpha = 255  # Default to fully opaque
                    if hasattr(element.stroke, "alpha") and element.stroke.alpha is not None:
                        stroke_alpha = element.stroke.alpha
                    # Apply element-level opacity if present
                    opacity_val = element.values.get('opacity')  # Read opacity attribute
                    if opacity_val is not None:
                        try:
                            opacity_float = float(opacity_val)
                            opacity_alpha = int(opacity_float * 255)
                            stroke_alpha = int((stroke_alpha / 255.0) * opacity_alpha)
                            logger.debug(f"Opacity applied to stroke: {opacity_val} -> alpha={stroke_alpha}")
                        except (ValueError, TypeError):
                            logger.debug(f"Could not parse opacity: {opacity_val}")
                    if stroke_alpha < 255:
                        PPTXNativeExporter._apply_transparency(
                            shape.line.color, stroke_alpha
                        )
                # Convert pixels to points (1 px = 0.75 pt at 96 DPI)
                shape.line.width = Pt(element.stroke_width * 0.75)
            else:
                shape.line.fill.background()

            return shape

        except Exception as e:
            logger.error(f"Path handling error: {e}")
            return

    @staticmethod

    def _handle_text(element, slide, tx, ty, svg_doc=None, svg_code=None):
        """
        Handle SVG Text elements and create PPTX text boxes.
        """
        try:
            # Extract text content
            text_content = ""
            if hasattr(element, "text"):
                text_content = str(element.text).strip()

            if not text_content:
                logger.debug("Text element has no content")
                return None

            # Get position and size
            x = (
                float(element.x)
                if hasattr(element, "x") and element.x is not None
                else 0
            )
            y = (
                float(element.y)
                if hasattr(element, "y") and element.y is not None
                else 0
            )

            # Get font size
            font_size = 12  # Default
            if hasattr(element, "font_size") and element.font_size is not None:
                try:
                    font_size = float(element.font_size)
                except (ValueError, TypeError):
                    font_size = 12

            # Get color
            color = element.fill if hasattr(element, "fill") else Color("black")
            rgb = PPTXNativeExporter._color_to_rgb(color, svg_doc, element.values if hasattr(element, 'values') else None, svg_code)
            if not rgb:
                rgb = RGBColor(0, 0, 0)  # Default to black

            # Create text box (estimate width based on text length)
            text_width = Inches(max(1.0, len(text_content) * 0.05))
            text_height = Inches(0.3)

            textbox = slide.shapes.add_textbox(tx(x), ty(y), text_width, text_height)
            text_frame = textbox.text_frame
            text_frame.word_wrap = True

            # Add text paragraph
            p = text_frame.paragraphs[0]
            p.text = text_content

            # Apply font properties
            run = p.runs[0] if p.runs else None
            if run:
                run.font.size = Pt(font_size)
                run.font.color.rgb = rgb
                logger.debug(
                    f"Text created: '{text_content}' at ({x}, {y}) with size {font_size}pt"
                )

            return textbox

        except Exception as e:
            logger.error(f"Text handling error: {e}")
            return None

    @staticmethod
    @staticmethod
    def _handle_image(element, slide, tx, ty, svg_doc=None, svg_code=None):
        """
        Handle SVG Image elements and create PPTX pictures.
        """
        try:
            # Get position
            x = (
                float(element.x)
                if hasattr(element, "x") and element.x is not None
                else 0
            )
            y = (
                float(element.y)
                if hasattr(element, "y") and element.y is not None
                else 0
            )

            # Get size
            width = (
                float(element.width)
                if hasattr(element, "width") and element.width is not None
                else 100
            )
            height = (
                float(element.height)
                if hasattr(element, "height") and element.height is not None
                else 100
            )

            # Get image href/data
            href = None
            if hasattr(element, "href") and element.href:
                href = element.href
            elif hasattr(element, "url") and element.url:
                href = element.url

            if not href:
                logger.debug("Image element has no href")
                return None

            # Handle data URIs or file paths
            if href.startswith("data:image"):
                # Data URI - need to decode and save temporarily
                logger.warning("Data URI images not fully supported yet")
                return None

            # For regular file paths, try to load the image
            try:
                picture = slide.shapes.add_picture(
                    href,
                    tx(x),
                    ty(y),
                    width=Inches(width * 0.01),
                    height=Inches(height * 0.01),
                )
                logger.debug(f"Image added: {href} at ({x}, {y})")
                return picture
            except Exception as img_err:
                logger.warning(f"Failed to add image from {href}: {img_err}")
                return None

        except Exception as e:
            logger.error(f"Image handling error: {e}")
            return None

    @staticmethod
    def generate_pptx_from_data(svg_code: str) -> bytes:
        """Generate PPTX from SVG code."""
        logger.info("PPTX generation started")
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        try:
            svg_doc = SVG.parse(io.StringIO(svg_code))
        except Exception as e:
            logger.error(f"SVG parse error: {e}")
            return b""

        # Get viewbox/dimensions
        vb = svg_doc.viewbox
        if vb and hasattr(vb, "width") and hasattr(vb, "height"):
            vw, vh = vb.width, vb.height
        else:
            vw = (
                float(svg_doc.width)
                if hasattr(svg_doc, "width") and svg_doc.width
                else 800
            )
            vh = (
                float(svg_doc.height)
                if hasattr(svg_doc, "height") and svg_doc.height
                else 600
            )

        if vw <= 0 or vh <= 0:
            vw, vh = 800, 600

        # Scale to fit PPTX slide (10 inches x 7.5 inches)
        scale_x = 10.0 / vw
        scale_y = 7.5 / vh

        def tx(v):
            return Inches(float(v) * scale_x)

        def ty(v):
            return Inches(float(v) * scale_y)

        # Process elements
        for element in svg_doc.elements():
            try:
                if isinstance(element, Text):
                    logger.debug(f"Processing Text element")
                    PPTXNativeExporter._handle_text(element, slide, tx, ty, svg_doc, svg_code)

                elif isinstance(element, Image):
                    logger.debug(f"Processing Image element")
                    PPTXNativeExporter._handle_image(element, slide, tx, ty, svg_doc, svg_code)

                elif isinstance(element, Shape):
                    logger.debug(f"Processing Shape element")
                    # Skip if both fill and stroke are None
                    if element.stroke.value is None and element.fill.value is None:
                        logger.debug("Skipping shape with no fill or stroke")
                        continue

                    PPTXNativeExporter._handle_path(element, slide, tx, ty, svg_doc, svg_code)

                elif isinstance(element, (Group, Use)):
                    logger.debug(
                        f"Skipping {type(element).__name__} element (not directly rendered)"
                    )
                    continue

            except Exception as e:
                logger.error(f"Element processing error: {e}", exc_info=True)
                continue

        # Save presentation
        pptx_stream = io.BytesIO()
        prs.save(pptx_stream)
        logger.info("PPTX generation successful")
        return pptx_stream.getvalue()
